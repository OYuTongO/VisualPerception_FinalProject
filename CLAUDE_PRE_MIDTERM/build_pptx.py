# -*- coding: utf-8 -*-
"""Build presentation_v3.pptx from HTML content"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colors ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0d, 0x11, 0x17)
SURFACE  = RGBColor(0x16, 0x1b, 0x22)
SURFACE2 = RGBColor(0x21, 0x26, 0x2d)
ACCENT   = RGBColor(0x58, 0xa6, 0xff)
ACCENT2  = RGBColor(0x3f, 0xb9, 0x50)
ACCENT3  = RGBColor(0xf7, 0x81, 0x66)
ACCENT4  = RGBColor(0xd2, 0xa8, 0xff)
TEXT     = RGBColor(0xe6, 0xed, 0xf3)
MUTED    = RGBColor(0x8b, 0x94, 0x9e)
BORDER   = RGBColor(0x30, 0x36, 0x3d)
WHITE    = RGBColor(0xff, 0xff, 0xff)

# ── Setup ────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

W = float(prs.slide_width)
H = float(prs.slide_height)

# ── Helpers ──────────────────────────────────────────────────────────────────
def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s

def tb(slide, text, l, t, w, h, sz=13, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, bg=None, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    if bg:
        box.fill.solid(); box.fill.fore_color.rgb = bg
    else:
        box.fill.background()
    box.line.fill.background()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size     = Pt(sz)
    r.font.bold     = bold
    r.font.italic   = italic
    r.font.color.rgb = color if color else TEXT
    return box

def box(slide, l, t, w, h, fill=None, border=None, bw=Pt(0.75)):
    from pptx.util import Emu
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border; shp.line.width = bw
    else:
        shp.line.fill.background()
    return shp

def section_label(s, text, t=Inches(0.55)):
    tb(s, text, Inches(0.75), t, Inches(8), Inches(0.28),
       sz=9, bold=True, color=ACCENT)

def h2(s, text, t=Inches(0.9)):
    tb(s, text, Inches(0.75), t, Inches(11.8), Inches(0.65),
       sz=26, bold=True, color=TEXT)

def card(slide, l, t, w, h, title=None, tcol=None, lines=None,
         accent=None, title_sz=12, line_sz=10):
    box(slide, l, t, w, h, fill=SURFACE, border=BORDER, bw=Pt(0.5))
    if accent:
        box(slide, l, t, Inches(0.045), h, fill=accent)
    xl = l + Inches(0.18)
    xw = w - Inches(0.28)
    y  = t + Inches(0.14)
    if title:
        tb(slide, title, xl, y, xw, Inches(0.32), sz=title_sz, bold=True,
           color=tcol or TEXT)
        y += Inches(0.32)
    if lines:
        for item in (lines if isinstance(lines, list) else [lines]):
            c = MUTED
            if isinstance(item, tuple):
                item, c = item
            tb(slide, item, xl + Inches(0.05), y, xw - Inches(0.05),
               Inches(0.26), sz=line_sz, color=c, wrap=True)
            y += Inches(0.28)

def pill(slide, text, l, t, color):
    """Small colored tag/pill"""
    box(slide, l, t, Inches(1.4), Inches(0.22), fill=SURFACE2, border=color, bw=Pt(0.4))
    tb(slide, text, l + Inches(0.1), t + Inches(0.02), Inches(1.2), Inches(0.2),
       sz=8, bold=True, color=color, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s1 = new_slide()
# accent bg strip
box(s1, 0, Inches(2.8), W, Inches(2.2), fill=RGBColor(0x0f, 0x16, 0x20))
tb(s1, "华东师范大学 · 视觉感知与前沿技术", 0, Inches(0.3), W, Inches(0.4),
   sz=11, color=ACCENT, align=PP_ALIGN.CENTER)
tb(s1, "🤟", 0, Inches(0.9), W, Inches(0.7), sz=32, align=PP_ALIGN.CENTER)
tb(s1, "基于手势识别的实时手语教学反馈系统", 0, Inches(1.75), W, Inches(0.75),
   sz=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s1, "ASL Real-time Teaching Feedback System via Gesture Recognition",
   0, Inches(2.55), W, Inches(0.4), sz=13, color=ACCENT, align=PP_ALIGN.CENTER)
tb(s1, "Gesture Recognition  ·  MediaPipe  ·  Real-time Feedback  ·  Teaching System",
   0, Inches(3.1), W, Inches(0.35), sz=12, color=MUTED, align=PP_ALIGN.CENTER)
# bottom bar
box(s1, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.45), fill=SURFACE, border=BORDER)
tb(s1, "📅  2026 / 05 / 21        👤  汇报人   陈卓然   潘南至",
   Inches(3.5), Inches(6.52), Inches(6.3), Inches(0.35),
   sz=11, color=MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ════════════════════════════════════════════════════════════════════════════
s2 = new_slide()
section_label(s2, "BACKGROUND & MOTIVATION")
h2(s2, "为什么做这个项目？")

# 3 stat cards
for i, (num, lbl) in enumerate([
    ("7,200万+", "全球聋哑人口\n依赖手语沟通"),
    ("26",       "ASL 字母手势\n（J、Z 为动态手势）"),
    ("∞",        "传统学习方式\n缺乏实时反馈"),
]):
    xl = Inches(0.75) + i * Inches(4.1)
    box(s2, xl, Inches(1.75), Inches(3.8), Inches(1.1), fill=SURFACE, border=BORDER)
    tb(s2, num, xl, Inches(1.82), Inches(3.8), Inches(0.55),
       sz=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    tb(s2, lbl, xl, Inches(2.35), Inches(3.8), Inches(0.45),
       sz=10, color=MUTED, align=PP_ALIGN.CENTER)

# 2 bottom cards
card(s2, Inches(0.75), Inches(3.1), Inches(5.8), Inches(2.1),
     title="⚠  现有问题", tcol=ACCENT3,
     lines=[
         "·  手语学习门槛高，缺乏即时纠错反馈",
         "·  传统教学依赖线下老师，成本高昂",
         "·  现有 App 功能单一，无闭环纠错机制",
     ])
card(s2, Inches(6.8), Inches(3.1), Inches(5.8), Inches(2.1),
     title="✓  我们的方向", tcol=ACCENT2,
     lines=[
         "·  工业相机实时识别，低延迟视觉反馈",
         "·  即时纠错反馈，拼写练习教学闭环",
         "·  纯 Python 实现，普通笔记本即可运行",
     ])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ════════════════════════════════════════════════════════════════════════════
s3 = new_slide()
section_label(s3, "SOLUTION OVERVIEW")
h2(s3, "整体思路：看手  →  识字母  →  即时反馈")

img_dir = os.path.dirname(os.path.abspath(__file__))
imgs = [
    ("gesture_before_detected.png", BORDER,   ACCENT,  "原始输入",
     "摄像头实时画面\n未做任何处理"),
    ("gesture_after_detected.png",  ACCENT,   ACCENT,  "MediaPipe 关键点检测",
     "21 个手部关键点提取\nBBox 归一化 → 42 维特征"),
    ("alpha.png",                   ACCENT2,  ACCENT2, "实时教学反馈",
     "识别结果驱动反馈角色\n即时正误提示 + 教学引导"),
]
for i, (fname, bcol, lcol, cap, desc) in enumerate(imgs):
    xl = Inches(0.75) + i * Inches(4.2)
    fpath = os.path.join(img_dir, fname)
    if os.path.exists(fpath):
        try:
            s3.shapes.add_picture(fpath, xl, Inches(1.85), Inches(3.85), Inches(2.5))
            box(s3, xl, Inches(1.85), Inches(3.85), Inches(2.5), border=bcol)
        except Exception:
            box(s3, xl, Inches(1.85), Inches(3.85), Inches(2.5), fill=SURFACE2, border=bcol)
            tb(s3, f"[ {fname} ]", xl, Inches(2.8), Inches(3.85), Inches(0.5),
               sz=9, color=MUTED, align=PP_ALIGN.CENTER)
    else:
        box(s3, xl, Inches(1.85), Inches(3.85), Inches(2.5), fill=SURFACE2, border=bcol)
        tb(s3, f"[ {fname} ]", xl, Inches(2.8), Inches(3.85), Inches(0.5),
           sz=9, color=MUTED, align=PP_ALIGN.CENTER)
    tb(s3, cap,  xl, Inches(4.5),  Inches(3.85), Inches(0.3), sz=11, bold=True, color=lcol, align=PP_ALIGN.CENTER)
    tb(s3, desc, xl, Inches(4.82), Inches(3.85), Inches(0.5), sz=9,  color=MUTED,  align=PP_ALIGN.CENTER)
    if i < 2:
        tb(s3, "→", Inches(4.65) + i * Inches(4.2), Inches(3.0),
           Inches(0.3), Inches(0.4), sz=18, color=ACCENT, align=PP_ALIGN.CENTER)

# tech tags row
tags = [("OpenCV","blue"),("MediaPipe Tasks API","green"),("RandomForest","purple"),
        ("Pygame","blue"),("LanXin Eagle Pro","green"),("Python 3.10+","red")]
tcols = {"blue":ACCENT,"green":ACCENT2,"purple":ACCENT4,"red":ACCENT3}
xpos = Inches(0.75)
for tag_text, tag_color in tags:
    tw = Inches(len(tag_text) * 0.12 + 0.4)
    box(s3, xpos, Inches(5.5), tw, Inches(0.28),
        fill=SURFACE2, border=tcols[tag_color], bw=Pt(0.4))
    tb(s3, tag_text, xpos + Inches(0.08), Inches(5.52), tw - Inches(0.12), Inches(0.24),
       sz=8, bold=True, color=tcols[tag_color])
    xpos += tw + Inches(0.12)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Highlights
# ════════════════════════════════════════════════════════════════════════════
s4 = new_slide()
section_label(s4, "HIGHLIGHTS")
h2(s4, "亮点细节")

# 4 metric cards
metrics = [
    ("57,339", "关键点样本\n(Kaggle 87k 提取)", ACCENT2),
    ("24",     "字母类别\n(排除 J、Z)", ACCENT2),
    ("98.70%", "RandomForest\n测试集准确率", ACCENT2),
    ("94.4%",  "最弱字母 N\n(与 M 混淆)", ACCENT3),
]
for i, (num, lbl, col) in enumerate(metrics):
    xl = Inches(0.75) + i * Inches(3.15)
    box(s4, xl, Inches(1.75), Inches(2.95), Inches(1.1), fill=SURFACE, border=BORDER)
    tb(s4, num, xl, Inches(1.82), Inches(2.95), Inches(0.5),
       sz=24, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s4, lbl, xl, Inches(2.32), Inches(2.95), Inches(0.5),
       sz=9, color=MUTED, align=PP_ALIGN.CENTER)

# 3 mode cards
modes = [
    ("📖  学习模式", "引导教学",  ACCENT2,
     ["Easy 词库取词，高亮目标字母",  "正确停留1.5秒自动推进，→跳过"]),
    ("⏱  测试模式", "评估测试",  ACCENT,
     ["每字母限时10秒，+10分/字母",   "弱点统计 + Top 5 排行榜"]),
    ("🕹  练习模式", "自由练习",  ACCENT4,
     ["全屏摄像头背景，骨架叠加",     "大字母 + 置信度 + 最近5字历史"]),
]
for i, (title, tag, col, lines) in enumerate(modes):
    xl = Inches(0.75) + i * Inches(4.15)
    box(s4, xl, Inches(3.1), Inches(3.95), Inches(1.45), fill=SURFACE, border=BORDER)
    box(s4, xl, Inches(3.1), Inches(3.95), Inches(0.22), fill=SURFACE2, border=BORDER)
    tb(s4, tag, xl + Inches(0.12), Inches(3.13), Inches(1.3), Inches(0.18),
       sz=8, bold=True, color=col)
    tb(s4, title, xl + Inches(0.12), Inches(3.35), Inches(3.7), Inches(0.3),
       sz=13, bold=True, color=TEXT)
    for j, line in enumerate(lines):
        tb(s4, "›  " + line, xl + Inches(0.18), Inches(3.7) + j * Inches(0.3),
           Inches(3.6), Inches(0.28), sz=10, color=MUTED)

# Mascot states row
states = [
    ("IDLE",    "等待 / 无手势", MUTED,   "站立待机  ·  轻微呼吸动画"),
    ("CORRECT", "识别正确 ✓",   ACCENT2, "举手欢呼  ·  绿色光晕"),
    ("WRONG",   "识别错误 ✗",   ACCENT3, "摊手摇头  ·  红色抖动"),
]
tb(s4, "反馈角色三态", Inches(0.75), Inches(4.75), Inches(2.5), Inches(0.28),
   sz=10, bold=True, color=ACCENT4)
for i, (state, sub, col, desc) in enumerate(states):
    xl = Inches(0.75) + i * Inches(3.0)
    box(s4, xl, Inches(5.1), Inches(2.8), Inches(1.0), fill=SURFACE2, border=col, bw=Pt(0.8))
    tb(s4, state, xl, Inches(5.2), Inches(2.8), Inches(0.3),
       sz=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s4, sub,   xl, Inches(5.5), Inches(2.8), Inches(0.25),
       sz=9,  color=col,  align=PP_ALIGN.CENTER)
    tb(s4, desc,  xl, Inches(5.77), Inches(2.8), Inches(0.25),
       sz=9,  color=MUTED, align=PP_ALIGN.CENTER)
# implementation note
card(s4, Inches(9.75), Inches(5.1), Inches(2.85), Inches(1.0),
     title="反馈角色实现", tcol=ACCENT4,
     lines=["纯 Pygame 绘制，零外部图片",
            "Lerp 插值过渡，约 0.3 秒",
            "mascot.set_state(...)"])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Technical Pipeline
# ════════════════════════════════════════════════════════════════════════════
s5 = new_slide()
section_label(s5, "TECHNICAL PIPELINE")
h2(s5, "端到端处理流程")

steps = [
    ("📷", "图像采集",     "LanXin Eagle Pro"),
    ("✋", "手部检测",     "MediaPipe HandLandmarker"),
    ("📐", "关键点归一化", "21 pts → 42-dim vector"),
    ("🌲", "分类推理",     "RandomForestClassifier"),
    ("🔁", "帧平滑",       "3-frame majority vote"),
    ("🎓", "教学逻辑",     "Pygame Scene Manager"),
]
sw = Inches(1.75)
for i, (icon, name, tech) in enumerate(steps):
    xl = Inches(0.75) + i * Inches(2.12)
    box(s5, xl, Inches(1.75), sw, Inches(1.1), fill=SURFACE, border=BORDER)
    tb(s5, icon, xl, Inches(1.8), sw, Inches(0.32), sz=16, align=PP_ALIGN.CENTER)
    tb(s5, name, xl, Inches(2.12), sw, Inches(0.28), sz=10, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    tb(s5, tech, xl, Inches(2.4),  sw, Inches(0.28), sz=8, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 5:
        tb(s5, "→", Inches(2.52) + i * Inches(2.12), Inches(2.1),
           Inches(0.3), Inches(0.3), sz=14, color=ACCENT, align=PP_ALIGN.CENTER)

# Two bottom cards
card(s5, Inches(0.75), Inches(3.2), Inches(5.8), Inches(2.8),
     title="技术选型依据", tcol=ACCENT,
     lines=[
         ("相机：LanXin Eagle Pro — 高帧率、低畸变、强弱光", TEXT),
         ("视觉：MediaPipe Tasks API (0.10.35+)", TEXT),
         ("模型：RandomForest — 速度快，无需 GPU", TEXT),
         ("界面：Pygame — 跨平台，帧级渲染控制", TEXT),
     ])
card(s5, Inches(6.8), Inches(3.2), Inches(5.8), Inches(2.8),
     title="关键设计决策", tcol=ACCENT4,
     lines=[
         "MediaPipe 0.10+ 移除 mp.solutions，迁移至 Tasks API",
         "坐标二次归一化到手部 BBox，准确率 92% → 98.70%",
         "排除 J、Z（动态手势），仅识别 24 个静态字母",
         "摄像头帧在子线程采集，主循环保持 30 FPS",
     ])

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Retrospective (3 points)
# ════════════════════════════════════════════════════════════════════════════
s6 = new_slide()
section_label(s6, "RETROSPECTIVE")
h2(s6, "难点复盘")

retros = [
    ("难点 01 · 特征被位置信息污染",
     "直接使用 MediaPipe 输出坐标，同一手势在不同位置 / 距离下特征向量差异极大",
     "MediaPipe 输出坐标是相对整张图片的归一化值，手偏左、偏右、靠近或远离摄像头，"
     "42 个特征数字完全不同。模型被迫同时学习手势形状与手的空间位置，特征空间被污染，泛化能力极差。",
     "✓  在 MediaPipe 之后增加 BBox 二次归一化：x_norm = (x − x_min) / w，"
     "彻底剔除位置与尺度因素。准确率从 ~92% 跳至 98.70%"),
    ("难点 02 · 逐帧输出与教学事件的根本矛盾",
     "静态分类器以 30 fps 持续输出字母，同一手势静止不动即可"刷穿"整个单词",
     "分类器每帧均产生结果，用户将手摆成字母 A 静止，系统就以 30 fps 速率持续输入 A，"
     "单词毫秒内拼完，教学节奏崩溃。本质矛盾：逐帧分类器输出连续信号，"
     "教学系统需要离散的手势事件。",
     "✓  引入 _last_confirmed + 3 帧 deque 缓冲，将分类器包装为状态机——"
     "3 帧全部一致且与上次不同时才触发一次事件，连续信号转化为离散手势事件流"),
    ("难点 03 · 单帧静态模型的固有局限",
     "模型只看单帧坐标，无时序信息，导致两类根本性问题无法被训练数据解决",
     "① 动态手势（J、Z）需要手部运动轨迹，单帧坐标无法区分，直接排除 2 个字母。"
     "② M（3 指折叠）与 N（4 指折叠）仅差一根手指微小弯曲，"
     "成为准确率最低的字母对（94.4% / 95.4%）。两者共同根源：模型缺乏时间维度感知。",
     "✓  现阶段：词库规避高频 MN 连续组合；后续：引入连续帧序列建模（光流 / LSTM），"
     "同步解决动态手势与相似手型问题"),
]
for i, (tag, title, body, resolve) in enumerate(retros):
    yt = Inches(1.75) + i * Inches(1.68)
    yh = Inches(1.6)
    box(s6, Inches(0.75),  yt, Inches(0.045), yh, fill=ACCENT3)
    box(s6, Inches(0.795), yt, Inches(11.8),  yh, fill=SURFACE, border=BORDER)
    tb(s6, tag,     Inches(0.95), yt + Inches(0.1),  Inches(11.4), Inches(0.22), sz=9,  bold=True, color=ACCENT3)
    tb(s6, title,   Inches(0.95), yt + Inches(0.32), Inches(11.4), Inches(0.28), sz=12, bold=True, color=TEXT)
    tb(s6, body,    Inches(0.95), yt + Inches(0.62), Inches(11.4), Inches(0.42), sz=10, color=MUTED, wrap=True)
    tb(s6, resolve, Inches(0.95), yt + Inches(1.06), Inches(11.4), Inches(0.36), sz=10, color=ACCENT2, wrap=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Future Enhancement
# ════════════════════════════════════════════════════════════════════════════
s7 = new_slide()
section_label(s7, "FUTURE ENHANCEMENT")
h2(s7, "后续深化改进方向")

enhancements = [
    ("01", "动态手势扩展（J & Z）",
     "LanXin Eagle Pro 高帧率（≥60fps）为时序建模创造条件，"
     "引入光流法或简单 LSTM（~15 帧序列），识别 J/Z 轨迹，补全全部 26 字母。",
     "连续帧关键点序列 → 光流 / LSTM → 轨迹分类"),
    ("02", "手势质量评分（超越对/错二分）",
     "以标准 ASL 关键点为模板，计算余弦相似度或关节角度偏差，"
     "输出百分制质量分，如[你的 A 大拇指弯度不足，得 83 分]。",
     "关键点模板匹配 → 余弦相似度 / 角度误差"),
    ("03", "个性化弱点追踪 + 自适应出题",
     "记录每个用户每字母历史错误率，生成 24 格字母热图；"
     "测试模式根据热图加权采样，类似 Anki 间隔重复算法。",
     "错误率热图 → 加权随机采样 → 个性化出题"),
    ("04", "增量学习 · 适应个体手型差异",
     "教学开始前 5 分钟校准，每字母比划 3~5 次，新样本加入训练集做增量微调，"
     "解决不同用户手的大小/比例/习惯差异。",
     "sklearn warm_start=True，RandomForest 增量添加树"),
    ("05", "双手识别 + 连续字母流输入（练习模式升级）",
     "MediaPipe 支持 num_hands=2，实现连续字母流自动分词（基于手势停顿时间窗口），"
     "让自由练习从"展示工具"升级为真正的"手语输入法"。",
     "双手并行 Recognizer → 停顿检测 → 自动分词"),
]
positions = [
    (Inches(0.75),  Inches(1.75), Inches(6.05), Inches(1.45)),
    (Inches(7.05),  Inches(1.75), Inches(6.05), Inches(1.45)),
    (Inches(0.75),  Inches(3.35), Inches(6.05), Inches(1.45)),
    (Inches(7.05),  Inches(3.35), Inches(6.05), Inches(1.45)),
    (Inches(0.75),  Inches(4.95), Inches(12.35), Inches(1.3)),
]
for (xl, yt, cw, ch), (num, title, body, tech) in zip(positions, enhancements):
    box(s7, xl,               yt, Inches(0.045), ch, fill=ACCENT)
    box(s7, xl + Inches(0.045), yt, cw - Inches(0.045), ch, fill=SURFACE, border=BORDER)
    tb(s7, num,   xl + Inches(0.15), yt + Inches(0.1),  Inches(0.45), Inches(0.45), sz=20, bold=True, color=RGBColor(0x58,0xa6,0xff))
    tb(s7, title, xl + Inches(0.65), yt + Inches(0.12), cw - Inches(0.8), Inches(0.28), sz=12, bold=True, color=TEXT)
    tb(s7, body,  xl + Inches(0.65), yt + Inches(0.42), cw - Inches(0.8), Inches(0.55), sz=10, color=MUTED, wrap=True)
    tb(s7, tech,  xl + Inches(0.65), yt + Inches(0.98) if ch > Inches(1.3) else yt + Inches(0.95),
       cw - Inches(0.8), Inches(0.22), sz=9, italic=True, color=ACCENT4, wrap=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Development Progress
# ════════════════════════════════════════════════════════════════════════════
s8 = new_slide()
section_label(s8, "DEVELOPMENT PROGRESS")
h2(s8, "开发进度总览")

phases = [
    ("Phase 0",   "项目初始化",                          1.0, ACCENT2, "✅ 完成"),
    ("Phase 1",   "数据集获取（Kaggle ASL Alphabet）",   1.0, ACCENT2, "✅ 完成"),
    ("Phase 2",   "MediaPipe 关键点提取 → landmarks.csv",1.0, ACCENT2, "✅ 完成"),
    ("Phase 3",   "模型训练 → 98.70% 准确率",            1.0, ACCENT2, "✅ 完成"),
    ("Phase 4",   "实时识别模块 Recognizer（摄像头 → 字母）", 1.0, ACCENT2, "✅ 完成"),
    ("Phase 5",   "Pygame 主框架 + 反馈角色 + 单词库",   0.3, ACCENT,  "🔨 进行中"),
    ("Phase 6–8", "学习 / 测试 / 练习三大模式",          0.0, BORDER,  "📋 待开始"),
    ("Phase 9–10","集成测试 · 优化 · 收尾文档",          0.0, BORDER,  "📋 待开始"),
]
for i, (phase, desc, prog, col, status) in enumerate(phases):
    yt = Inches(1.75) + i * Inches(0.6)
    is_active = "进行中" in status
    bg_col = RGBColor(0x10, 0x18, 0x28) if is_active else SURFACE
    bd_col = ACCENT if is_active else BORDER
    box(s8, Inches(0.75), yt, Inches(11.85), Inches(0.52), fill=bg_col, border=bd_col)
    tb(s8, phase, Inches(0.95), yt + Inches(0.12), Inches(1.0), Inches(0.3),
       sz=10, bold=True, color=ACCENT if is_active else MUTED)
    tb(s8, desc, Inches(2.1), yt + Inches(0.12), Inches(7.5), Inches(0.3),
       sz=12, bold=is_active, color=ACCENT if is_active else TEXT)
    # progress bar bg
    box(s8, Inches(9.85), yt + Inches(0.18), Inches(1.8), Inches(0.16), fill=SURFACE2, border=None)
    if prog > 0:
        box(s8, Inches(9.85), yt + Inches(0.18), Inches(1.8 * prog), Inches(0.16), fill=col, border=None)
    tb(s8, status, Inches(11.7), yt + Inches(0.12), Inches(0.9), Inches(0.3),
       sz=10, color=ACCENT if is_active else (ACCENT2 if "完成" in status else MUTED))

# tags
for j, (tag, col) in enumerate([("CV 核心已完成 ✓", ACCENT2),
                                  ("模型精度达标 ✓",  ACCENT),
                                  ("教学反馈层构建中", ACCENT3)]):
    xl = Inches(0.75) + j * Inches(2.8)
    box(s8, xl, Inches(6.6), Inches(2.55), Inches(0.32), fill=SURFACE2, border=col, bw=Pt(0.5))
    tb(s8, tag, xl + Inches(0.1), Inches(6.63), Inches(2.35), Inches(0.26),
       sz=10, bold=True, color=col)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Roadmap
# ════════════════════════════════════════════════════════════════════════════
s9 = new_slide()
section_label(s9, "ROADMAP")
h2(s9, "后续计划 · 期末目标")

# Timeline
timeline = [
    ("Phase 0–4", "已完成\nCV 管线全通",  ACCENT2),
    ("Phase 5",   "进行中\n主框架+反馈",  ACCENT),
    ("Phase 6",   "学习模式\n引导教学",   MUTED),
    ("Phase 7",   "测试模式\n评估计分",   MUTED),
    ("Phase 8",   "练习模式\n实时展示",   MUTED),
    ("Phase 9–10","期末交付\nDemo+文档",  MUTED),
]
# draw line
box(s9, Inches(1.5), Inches(2.32), Inches(10.4), Inches(0.04), fill=BORDER)
for i, (phase, label, col) in enumerate(timeline):
    xl = Inches(1.5) + i * Inches(2.08)
    is_done = col == ACCENT2
    is_now  = col == ACCENT
    dot_fill = col if (is_done or is_now) else BG
    dot_size = Inches(0.22)
    box(s9, xl - dot_size/2, Inches(2.32) - dot_size/2 + Inches(0.02),
        dot_size, dot_size, fill=dot_fill, border=col, bw=Pt(1.5))
    tb(s9, phase, xl - Inches(0.8), Inches(1.65), Inches(1.6), Inches(0.25),
       sz=9, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s9, label, xl - Inches(0.7), Inches(2.65), Inches(1.4), Inches(0.45),
       sz=9, color=col, align=PP_ALIGN.CENTER)

# 3 goal cards
goals = [
    ("🎯", "近期目标（2 周内）",
     "完成 SceneManager、反馈角色三态动画、单词库\n实现 Menu → LearnMode 可运行教学闭环"),
    ("⚡", "性能目标",
     "识别延迟 < 200ms  ·  30 FPS 稳定运行\n普通笔记本（无 GPU）流畅跑全流程"),
    ("🏁", "期末交付物",
     "三模式完整教学系统  ·  排行榜\n演示视频 ≥ 2 分钟  ·  GitHub tag v1.0"),
]
for i, (icon, title, desc) in enumerate(goals):
    xl = Inches(0.75) + i * Inches(4.2)
    box(s9, xl, Inches(3.3), Inches(3.95), Inches(1.6), fill=SURFACE, border=BORDER)
    tb(s9, icon,  xl, Inches(3.4),  Inches(3.95), Inches(0.4), sz=20, align=PP_ALIGN.CENTER)
    tb(s9, title, xl, Inches(3.82), Inches(3.95), Inches(0.3), sz=12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    tb(s9, desc,  xl, Inches(4.15), Inches(3.95), Inches(0.55), sz=10, color=MUTED, align=PP_ALIGN.CENTER)

# core innovations bar
box(s9, Inches(0.75), Inches(5.2), Inches(11.85), Inches(0.6), fill=SURFACE, border=BORDER)
tb(s9, "🤟  核心创新：", Inches(0.9), Inches(5.28), Inches(2.0), Inches(0.28),
   sz=11, bold=True, color=TEXT)
tb(s9, "LanXin Eagle Pro 工业相机  ·  MediaPipe Tasks API  ·  BBox 二次归一化  ·  3 帧平滑触发  ·  纯 Pygame 反馈角色  ·  单词拼写教学反馈闭环",
   Inches(2.85), Inches(5.3), Inches(9.6), Inches(0.42), sz=10, color=MUTED, wrap=True)

# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(img_dir, "presentation_v3.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
