"""
Convert presentation_new.html to PPTX by screenshotting each slide with Playwright.
"""
import os
import sys
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io

HTML_FILE = Path(__file__).parent / "presentation_new.html"
OUTPUT    = Path(__file__).parent / "presentation_new.pptx"
TOTAL_SLIDES = 10
SLIDE_W = 1920
SLIDE_H = 1080

def capture_slides():
    screenshots = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H})
        page.goto(HTML_FILE.as_uri())
        page.wait_for_timeout(800)  # let fonts / animations settle

        for i in range(1, TOTAL_SLIDES + 1):
            # activate the target slide via JS
            page.evaluate(f"goTo({i})")
            page.wait_for_timeout(600)
            png = page.screenshot(full_page=False)
            screenshots.append(png)
            print(f"  captured slide {i}/{TOTAL_SLIDES}")

        browser.close()
    return screenshots

def build_pptx(screenshots):
    prs = Presentation()
    # widescreen 16:9
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for idx, png_bytes in enumerate(screenshots):
        slide = prs.slides.add_slide(blank_layout)
        img_stream = io.BytesIO(png_bytes)
        slide.shapes.add_picture(
            img_stream,
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        print(f"  added slide {idx+1}/{len(screenshots)} to PPTX")

    prs.save(OUTPUT)
    print(f"\nSaved: {OUTPUT}")

if __name__ == "__main__":
    print("Step 1: capturing slides with Playwright …")
    shots = capture_slides()
    print("Step 2: building PPTX …")
    build_pptx(shots)
